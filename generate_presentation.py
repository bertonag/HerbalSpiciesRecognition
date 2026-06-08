"""
Generates a PowerPoint presentation for HerbScan.

Metrics are loaded dynamically from runs/train/metrics.json (written by
evaluate.py).  Falls back to hardcoded values if the file is absent.

Workflow:
    python train.py
    python evaluate.py              # writes runs/train/metrics.json
    python generate_presentation.py # picks up fresh metrics
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE_DIR     = Path(__file__).parent
OUT          = BASE_DIR / "HerbScan_Presentation.pptx"
METRICS_JSON = BASE_DIR / "runs" / "train" / "metrics.json"

# ── hardcoded fallback ────────────────────────────────────────────────────────
_FALLBACK = {
    "overall": {
        "mAP50": 0.465, "mAP50_95": 0.313,
        "precision": 0.451, "recall": 0.480, "F1": 0.465,
    },
    "per_class": [
        {"class": "Bitter Leaf",       "P": 0.913, "R": 0.500, "F1": 0.646, "mAP50": 0.649, "mAP50_95": 0.385},
        {"class": "False Daisy",       "P": 1.000, "R": 0.975, "F1": 0.987, "mAP50": 0.995, "mAP50_95": 0.611},
        {"class": "Momordica foetida", "P": 0.321, "R": 0.534, "F1": 0.401, "mAP50": 0.535, "mAP50_95": 0.463},
        {"class": "Perilla frutescens","P": 0.539, "R": 1.000, "F1": 0.700, "mAP50": 0.995, "mAP50_95": 0.796},
        {"class": "Plectranthus pros.","P": 0.381, "R": 0.250, "F1": 0.302, "mAP50": 0.148, "mAP50_95": 0.048},
        {"class": "Plectrarithus cya.","P": 0.000, "R": 0.000, "F1": 0.000, "mAP50": 0.020, "mAP50_95": 0.006},
        {"class": "Sugar cane",        "P": 0.152, "R": 0.100, "F1": 0.121, "mAP50": 0.051, "mAP50_95": 0.034},
        {"class": "Peppermint",        "P": 0.298, "R": 0.479, "F1": 0.367, "mAP50": 0.331, "mAP50_95": 0.164},
    ],
    "converged_epoch": 88,
}


def load_metrics() -> dict:
    if METRICS_JSON.exists():
        with open(METRICS_JSON) as f:
            data = json.load(f)
        print(f"[INFO] Loaded metrics from {METRICS_JSON}")
        return data
    print(f"[WARN] {METRICS_JSON} not found — using fallback metrics.")
    return _FALLBACK

# ── colour palette ────────────────────────────────────────────────────────────
GREEN_DARK   = RGBColor(0x1B, 0x5E, 0x20)   # deep forest green
GREEN_MID    = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT  = RGBColor(0xA5, 0xD6, 0xA7)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xF1, 0xF8, 0xE9)
DARK_GREY    = RGBColor(0x37, 0x47, 0x37)
ACCENT       = RGBColor(0xFF, 0xA0, 0x00)   # amber accent

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width_pt=0):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=16, color=DARK_GREY,
                   bullet_char="▸ ", bold_first=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if bold_first and p == tf.paragraphs[0]:
            run.font.bold = True
    return txb


def header_bar(slide, title, subtitle=None):
    """Green header bar at top of slide."""
    add_rect(slide, 0, 0, 13.33, 1.2, GREEN_DARK)
    add_text_box(slide, title, 0.3, 0.1, 12.5, 0.7,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.75, 12.5, 0.4,
                     font_size=14, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)
    # thin accent line
    add_rect(slide, 0, 1.2, 13.33, 0.05, ACCENT)


def footer_bar(slide, note=""):
    add_rect(slide, 0, 7.1, 13.33, 0.4, GREEN_DARK)
    add_text_box(slide, "HerbScan | Medicinal Herb Recognition | 2026",
                 0.2, 7.12, 8, 0.3,
                 font_size=10, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)
    if note:
        add_text_box(slide, note, 8.5, 7.12, 4.6, 0.3,
                     font_size=10, color=ACCENT, align=PP_ALIGN.RIGHT)


def section_card(slide, l, t, w, h, title, body_lines, title_bg=GREEN_MID):
    add_rect(slide, l, t, w, 0.45, title_bg)
    add_text_box(slide, title, l+0.1, t+0.04, w-0.2, 0.38,
                 font_size=14, bold=True, color=WHITE)
    add_rect(slide, l, t+0.45, w, h-0.45, LIGHT_GREY)
    y_off = t + 0.55
    for line in body_lines:
        add_text_box(slide, "• " + line, l+0.15, y_off, w-0.25, 0.38,
                     font_size=12, color=DARK_GREY)
        y_off += 0.38


def metric_box(slide, l, t, w, h, value, label, bg=GREEN_MID):
    add_rect(slide, l, t, w, h, bg)
    add_text_box(slide, value, l, t+0.1, w, h*0.5,
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l, t+h*0.55, w, h*0.4,
                 font_size=11, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)


def build():
    m           = load_metrics()
    ov          = m["overall"]
    per_class   = m["per_class"]
    conv_epoch  = m.get("converged_epoch") or 88

    sorted_cls  = sorted(per_class, key=lambda c: c["mAP50"], reverse=True)
    best_cls    = sorted_cls[0]  if sorted_cls else {"class": "—", "mAP50": 0}
    worst_cls   = sorted_cls[-1] if sorted_cls else {"class": "—", "mAP50": 0}
    second_best = sorted_cls[1]  if len(sorted_cls) > 1 else {"class": "—", "mAP50": 0}

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, GREEN_DARK)
    # decorative green gradient blocks
    add_rect(sl, 0, 5.5, 13.33, 2.0, GREEN_MID)
    add_rect(sl, 0, 6.8, 13.33, 0.7, rgb(0x10, 0x40, 0x10))
    # leaf accent
    add_rect(sl, 11.5, 0, 1.83, 7.5, rgb(0x1A, 0x69, 0x1A))

    add_text_box(sl, "🌿 HerbScan", 0.6, 0.8, 10.5, 1.1,
                 font_size=52, bold=True, color=WHITE)
    add_text_box(sl,
        "Medicinal Herbal Plant Recognition\nUsing YOLOv8 Instance Segmentation",
        0.6, 1.9, 10.5, 1.4,
        font_size=26, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)
    add_text_box(sl,
        "With Continuous Learning & Expert Consultation",
        0.6, 3.25, 10.5, 0.6,
        font_size=18, italic=True, color=ACCENT)
    add_text_box(sl, "Gilbert Nyakana", 0.6, 4.2, 8, 0.5,
                 font_size=18, color=WHITE)
    add_text_box(sl, "Makerere University, Kampala, Uganda | 2026",
                 0.6, 4.7, 8, 0.4, font_size=14, color=GREEN_LIGHT)
    add_text_box(sl,
        "GitHub: https://github.com/bertonag/HerbalSpiciesRecognition",
        0.6, 5.6, 12, 0.5, font_size=14, color=ACCENT)
    add_text_box(sl, "Slide 1 / 14", 0.6, 6.9, 12, 0.4,
                 font_size=11, color=GREEN_LIGHT, align=PP_ALIGN.RIGHT)

    # ── SLIDE 2 — MOTIVATION ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Motivation", "Why does herbal plant recognition matter?")
    footer_bar(sl, "2 / 14")

    section_card(sl, 0.3, 1.4, 5.9, 2.5, "The Problem", [
        "60%+ of Uganda's population relies on herbal medicine",
        "Misidentification causes poisonings & treatment failures",
        "Expert botanists are scarce and geographically concentrated",
        "No low-cost automated identification tool exists for the region",
    ])
    section_card(sl, 6.4, 1.4, 6.6, 2.5, "Our Goal", [
        "Build a real-time herbal plant recognition system",
        "Run on low-cost hardware (CPU-only)",
        "Provide unknown-herb escalation to human experts",
        "Continuously improve via user corrections",
    ])

    add_rect(sl, 0.3, 4.1, 12.7, 0.05, GREEN_LIGHT)
    add_text_box(sl,
        '"An automated, accessible, and continuously improving system for '
        'identifying medicinal herbs at the point of need."',
        0.5, 4.2, 12.3, 0.9,
        font_size=16, italic=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    for i, (icon, label) in enumerate([
        ("🌿", "16 Herb Species"), ("📷", "Camera Capture"),
        ("🤖", "YOLOv8 AI"), ("♻️", "Continuous Learning")
    ]):
        x = 0.8 + i * 3.1
        add_rect(sl, x, 5.3, 2.5, 1.5, GREEN_LIGHT)
        add_text_box(sl, icon, x, 5.35, 2.5, 0.7,
                     font_size=28, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, x, 5.9, 2.5, 0.5,
                     font_size=13, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # ── SLIDE 3 — DATASET ────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Dataset", "Herbal Plants SpeciesInstSeg (Roboflow)")
    footer_bar(sl, "3 / 14")

    # Split stats
    for i, (n, label, bg) in enumerate([
        ("57", "Train Images", GREEN_DARK),
        ("10", "Val Images",   GREEN_MID),
        ("7",  "Test Images",  rgb(0x55, 0x8B, 0x2F)),
        ("307","Annotations",  ACCENT),
    ]):
        x = 0.3 + i * 3.2
        metric_box(sl, x, 1.4, 2.9, 1.6, n, label, bg)

    add_text_box(sl, "16 Medicinal Herb Classes:", 0.3, 3.2, 12.5, 0.45,
                 font_size=15, bold=True, color=GREEN_DARK)

    herbs = [
        "Bitter Leaf (Mululuza)", "Boerhavia diffusa", "Ceratonia siliqua",
        "Chenopodium album", "False Daisy (Mutaayiza)", "Himalayan Balsam",
        "Hoslundia opposita", "Justicia pectoralis", "Leucaena leucocephala",
        "Momordica foetida", "Perilla frutescens", "Plectranthus prostratus",
        "Plectrarithus cyaneus", "Sugar cane", "Peppermint", "New (unknown)",
    ]
    cols = 4
    for idx, herb in enumerate(herbs):
        col = idx % cols
        row = idx // cols
        x = 0.3 + col * 3.25
        y = 3.7 + row * 0.42
        add_rect(sl, x, y, 3.1, 0.38, LIGHT_GREY)
        add_text_box(sl, f"• {herb}", x+0.1, y+0.03, 3.0, 0.35,
                     font_size=11, color=DARK_GREY)

    add_text_box(sl,
        "⚠  Mixed annotations: 285 polygons + 22 bounding boxes → converted to polygons",
        0.3, 7.0, 12.5, 0.35, font_size=11, italic=True, color=rgb(0xB7, 0x50, 0x00))

    # ── SLIDE 4 — ANNOTATION PRE-PROCESSING ──────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Annotation Pre-processing",
               "Converting bounding boxes to polygons for YOLOv8-seg")
    footer_bar(sl, "4 / 14")

    add_text_box(sl, "The Challenge:", 0.3, 1.4, 12.5, 0.45,
                 font_size=16, bold=True, color=GREEN_DARK)
    add_text_box(sl,
        "YOLOv8-seg crashes with a RuntimeError when a mini-batch contains "
        "bounding-box only labels — the segmentation loss requires mask data.",
        0.3, 1.85, 12.5, 0.55, font_size=14, color=DARK_GREY)

    add_rect(sl, 0.3, 2.55, 5.9, 2.8, LIGHT_GREY)
    add_rect(sl, 0.3, 2.55, 5.9, 0.4, GREEN_MID)
    add_text_box(sl, "Before (5 tokens — bbox format)", 0.4, 2.58, 5.7, 0.35,
                 font_size=13, bold=True, color=WHITE)
    add_text_box(sl, "class_id  cx  cy  w  h",
                 0.5, 3.05, 5.6, 0.5, font_size=15, color=GREEN_DARK,
                 align=PP_ALIGN.CENTER)
    add_text_box(sl, "0  0.512  0.488  0.320  0.415",
                 0.5, 3.55, 5.6, 0.45, font_size=13, color=DARK_GREY,
                 align=PP_ALIGN.CENTER)
    add_text_box(sl, "↓ Not compatible with segmentation head",
                 0.5, 4.0, 5.6, 0.4, font_size=12, italic=True,
                 color=rgb(0xC6, 0x28, 0x28), align=PP_ALIGN.CENTER)

    # arrow
    add_text_box(sl, "→", 6.35, 3.7, 0.7, 0.6, font_size=32,
                 bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_rect(sl, 7.2, 2.55, 5.8, 2.8, LIGHT_GREY)
    add_rect(sl, 7.2, 2.55, 5.8, 0.4, GREEN_DARK)
    add_text_box(sl, "After (9 tokens — polygon format)", 7.3, 2.58, 5.6, 0.35,
                 font_size=13, bold=True, color=WHITE)
    add_text_box(sl, "class_id  x₁ y₁  x₂ y₂  x₃ y₃  x₄ y₄",
                 7.3, 3.05, 5.6, 0.5, font_size=13, color=GREEN_DARK,
                 align=PP_ALIGN.CENTER)
    add_text_box(sl, "4 corner points (clockwise)\nclamped to [0, 1]",
                 7.3, 3.55, 5.6, 0.6, font_size=12, italic=True,
                 color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_text_box(sl, "✓ Works with segmentation head",
                 7.3, 4.15, 5.6, 0.4, font_size=12,
                 color=rgb(0x1B, 0x5E, 0x20), align=PP_ALIGN.CENTER)

    add_text_box(sl,
        "Result: All 307 annotations retained  |  "
        "No training data discarded  |  "
        "convert_all_labels() in train.py",
        0.3, 5.6, 12.5, 0.5,
        font_size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # ── SLIDE 5 — MODEL ARCHITECTURE ─────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Model Architecture", "YOLOv8n Instance Segmentation")
    footer_bar(sl, "5 / 14")

    # Architecture blocks
    blocks = [
        ("Input\n640×640", 0.3, GREEN_MID),
        ("CSP-Darknet\nBackbone", 2.4, GREEN_DARK),
        ("PANet\nNeck", 4.5, GREEN_MID),
        ("Detection\nHead ×3", 6.6, GREEN_DARK),
        ("Seg Head\n32 protos", 8.7, GREEN_MID),
        ("Output\nMasks+BBoxes", 10.8, rgb(0x33, 0x69, 0x1E)),
    ]
    for label, x, bg in blocks:
        add_rect(sl, x, 1.5, 1.9, 1.2, bg)
        add_text_box(sl, label, x, 1.55, 1.9, 1.1,
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrows
    for x in [2.2, 4.3, 6.4, 8.5, 10.7]:
        add_text_box(sl, "→", x, 1.85, 0.25, 0.5, font_size=16,
                     bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

    specs = [
        ("Parameters",   "3.26M",   0.3),
        ("GFLOPs",       "11.4",    2.8),
        ("COCO mAP50",   "37.3",    5.3),
        ("Inference",    "~89ms",   7.8),
        ("Classes",      "16",      10.3),
    ]
    for label, val, x in specs:
        add_rect(sl, x, 3.0, 2.3, 1.1, LIGHT_GREY)
        add_text_box(sl, val, x, 3.0, 2.3, 0.6,
                     font_size=24, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, x, 3.55, 2.3, 0.45,
                     font_size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "Training Configuration:", 0.3, 4.3, 12.5, 0.45,
                 font_size=15, bold=True, color=GREEN_DARK)
    train_items = [
        "Epochs: 150 (early stop patience=40)   |   Converged: epoch 88",
        "Optimiser: AdamW  |  lr₀=0.01  |  Cosine annealing",
        "Augmentation: Mosaic p=1.0, MixUp α=0.15, Copy-Paste p=0.2",
        "Geometric: rotation ±15°, scale ×[0.4–1.6], shear ±5°, flips",
    ]
    for i, item in enumerate(train_items):
        add_text_box(sl, "• " + item, 0.4, 4.8 + i*0.38, 12.4, 0.38,
                     font_size=12, color=DARK_GREY)

    # ── SLIDE 6 — SYSTEM OVERVIEW ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "System Overview", "End-to-end pipeline")
    footer_bar(sl, "6 / 14")

    # Pipeline diagram
    pipeline = [
        ("📷 Camera\nor Upload", GREEN_DARK,   0.2),
        ("🔍 YOLOv8\nInference",  GREEN_MID,    2.4),
        ("📊 Results\nDisplay",   GREEN_DARK,   4.6),
        ("❓ Unknown?\nExpert Form",rgb(0x4A,0x6A,0x1A), 6.8),
        ("📝 Feedback\nCorrection",ACCENT,       9.0),
        ("♻️ Background\nFine-tune",GREEN_DARK, 11.2),
    ]
    for label, bg, x in pipeline:
        add_rect(sl, x, 1.5, 1.9, 1.6, bg)
        add_text_box(sl, label, x, 1.55, 1.9, 1.5,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for x in [2.1, 4.3, 6.5, 8.7, 11.0]:
        add_text_box(sl, "→", x, 2.1, 0.35, 0.5, font_size=18,
                     bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

    # Feature callout boxes
    features = [
        ("Per-Class\nDeduplication",
         "Each herb shown\nonce with best\nconfidence + count", 0.3, 3.4),
        ("Unknown Herb\nDetection",
         "Confidence <0.35\nor 'New' class\n→ expert form", 3.6, 3.4),
        ("Expert\nConsultation",
         "Submit description,\nlocation, contact\nfor expert review", 6.9, 3.4),
        ("Continuous\nLearning",
         "Corrections auto-\ntrigger background\nfine-tuning", 10.2, 3.4),
    ]
    for title, body, x, y in features:
        add_rect(sl, x, y, 2.95, 0.45, GREEN_MID)
        add_text_box(sl, title, x+0.05, y+0.03, 2.85, 0.4,
                     font_size=12, bold=True, color=WHITE)
        add_rect(sl, x, y+0.45, 2.95, 1.6, LIGHT_GREY)
        add_text_box(sl, body, x+0.1, y+0.55, 2.75, 1.4,
                     font_size=11, color=DARK_GREY)

    add_text_box(sl, "Low lr fine-tuning (lr=0.0005) prevents catastrophic forgetting",
                 0.3, 6.7, 12.7, 0.4,
                 font_size=12, italic=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # ── SLIDE 7 — WEB INTERFACE ───────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Web Interface", "Gradio 5.x – Browser-based GUI")
    footer_bar(sl, "7 / 14")

    panels = [
        ("📷 Input Panel",
         ["Camera capture via browser MediaDevices API",
          "Image upload (JPEG/PNG/BMP)",
          "Detect button triggers YOLOv8 inference",
          "Annotated result image with masks"],
         0.3, 1.4, 5.9),
        ("📊 Results Panel",
         ["Per-class deduplication: each herb shown once",
          "Best confidence score displayed",
          "Region count (how many instances found)",
          "Status: Identified or Unknown"],
         6.5, 1.4, 6.5),
        ("❓ Expert Form",
         ["Shown only when unknown herb detected",
          "User submits description + location",
          "Unique reference ID issued (HRB-XXXX)",
          "Logged to expert_queries.json"],
         0.3, 4.1, 5.9),
        ("✏️ Feedback Loop",
         ["Select wrong prediction from dropdown",
          "Choose correct herb species",
          "Submit triggers background fine-tuning",
          "Model hot-reloaded after training"],
         6.5, 4.1, 6.5),
    ]
    for title, items, x, y, w in panels:
        add_rect(sl, x, y, w, 0.42, GREEN_DARK)
        add_text_box(sl, title, x+0.1, y+0.04, w-0.15, 0.35,
                     font_size=13, bold=True, color=WHITE)
        for i, item in enumerate(items):
            add_text_box(sl, f"• {item}", x+0.15, y+0.5+i*0.35, w-0.2, 0.35,
                         font_size=11, color=DARK_GREY)

    # ── SLIDE 8 — DISCRIMINATIVE VISUAL PROPERTIES ───────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Discriminative Visual Properties",
               "What makes each herb visually unique?")
    footer_bar(sl, "8 / 14")

    species_props = [
        ("Perilla frutescens",    "🎨 Colour",   "Purple-red abaxial surface\n(H=270–320°, a*>0 in LAB)\nUnique in dataset"),
        ("False Daisy",           "🔲 Texture",  "Deeply serrated margins\nDense appressed hairs\nHigh LBP transition rate"),
        ("Sugar cane",            "📐 Shape",    "Aspect ratio ≈ 8:1\nParallel veins\nHighly elongated"),
        ("Momordica foetida",     "📐 Shape",    "Palmate lobed leaf\nSolidity ≈ 0.6\nLow circularity"),
        ("Bitter Leaf",           "🔲 Texture",  "Green-on-green camouflage\nDiscriminated by venation\nGLCM + Gabor needed"),
        ("Peppermint",            "📐+🎨",       "Square stems + opposite leaves\nSilvery surface sheen\nDistinctive aspect ratio"),
    ]
    cols_per_row = 3
    for idx, (species, feat_type, desc) in enumerate(species_props):
        col = idx % cols_per_row
        row = idx // cols_per_row
        x = 0.25 + col * 4.35
        y = 1.45 + row * 2.55
        feat_bg = (GREEN_MID if "Colour" in feat_type
                   else rgb(0x4A, 0x6A, 0x1A) if "Texture" in feat_type
                   else ACCENT)
        add_rect(sl, x, y, 4.1, 0.4, feat_bg)
        add_text_box(sl, f"{feat_type}  —  {species}", x+0.1, y+0.04, 3.9, 0.34,
                     font_size=11, bold=True, color=WHITE)
        add_rect(sl, x, y+0.4, 4.1, 1.9, LIGHT_GREY)
        for li, line in enumerate(desc.split("\n")):
            add_text_box(sl, f"• {line}", x+0.15, y+0.52+li*0.42, 3.8, 0.4,
                         font_size=11, color=DARK_GREY)

    add_text_box(sl,
        "Key insight: species with unique colour OR extreme shape are learnable from few images. "
        "Green-on-green species need texture features AND more training data.",
        0.3, 6.65, 12.7, 0.42, font_size=11, italic=True,
        color=GREEN_DARK, align=PP_ALIGN.CENTER)

    # ── SLIDE 9 — COLOUR & TEXTURE FEATURES ──────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Feature Extraction: Colour & Texture",
               "What to extract — and why it matters for herbs")
    footer_bar(sl, "9 / 14")

    # Colour column
    add_rect(sl, 0.25, 1.35, 6.1, 0.4, GREEN_DARK)
    add_text_box(sl, "🎨  COLOUR FEATURES", 0.35, 1.38, 5.9, 0.35,
                 font_size=14, bold=True, color=WHITE)
    colour_items = [
        ("HSV Histograms",
         "Hue channel: Perilla = H≈300° (purple), all others H≈90–150° (green)\n"
         "32 bins/channel → 96-D descriptor. Saturation detects dried vs. fresh."),
        ("LAB Colour Space",
         "a* axis (red–green): large positive value uniquely for Perilla.\n"
         "L* captures canopy-shade lighting variation. Perceptually uniform."),
        ("Vegetation Indices",
         "ExG = 2G−R−B  segments plant from background.\n"
         "VARI = (G−R)/(G+R−B)  robust to illumination variation.\n"
         "ExR = 1.4R−G  amplifies anthocyanin red in Perilla."),
        ("Chromaticity r,g",
         "r=R/(R+G+B), g=G/(R+G+B) — illumination-invariant ratios.\n"
         "Stable across exposure and white-balance variation in field photos."),
    ]
    y0 = 1.85
    for title, body in colour_items:
        add_rect(sl, 0.25, y0, 6.1, 0.3, rgb(0xE8, 0xF5, 0xE9))
        add_text_box(sl, title, 0.35, y0+0.02, 5.9, 0.27,
                     font_size=11, bold=True, color=GREEN_DARK)
        y0 += 0.3
        for line in body.split("\n"):
            add_text_box(sl, f"  {line}", 0.35, y0, 5.9, 0.28,
                         font_size=10, color=DARK_GREY)
            y0 += 0.28
        y0 += 0.04

    # Texture column
    add_rect(sl, 6.98, 1.35, 6.1, 0.4, rgb(0x4A, 0x6A, 0x1A))
    add_text_box(sl, "🔲  TEXTURE FEATURES", 7.08, 1.38, 5.9, 0.35,
                 font_size=14, bold=True, color=WHITE)
    texture_items = [
        ("Gabor Filter Banks  (6θ × 4σ)",
         "Captures venation, surface hairiness, and margin edge patterns.\n"
         "Multi-scale (4 scales) robust to varying camera distances.\n"
         "Response: mean + variance per filter → 48-D descriptor."),
        ("Local Binary Patterns (LBP)",
         "Rotation-invariant micro-texture. Multi-radius (r=1,2,3).\n"
         "Smooth waxy (Perilla) → few transitions.\n"
         "Hairy surface (False Daisy) → many transitions."),
        ("GLCM — Haralick Features",
         "5 key statistics: Energy, Contrast, Correlation,\n"
         "Homogeneity, Entropy. 4 orientations → rotation invariant.\n"
         "Contrast high for deeply veined leaves; Energy high for smooth."),
    ]
    y0 = 1.85
    for title, body in texture_items:
        add_rect(sl, 6.98, y0, 6.1, 0.3, rgb(0xF1, 0xF8, 0xE9))
        add_text_box(sl, title, 7.08, y0+0.02, 5.9, 0.27,
                     font_size=11, bold=True, color=rgb(0x2E, 0x4A, 0x1A))
        y0 += 0.3
        for line in body.split("\n"):
            add_text_box(sl, f"  {line}", 7.08, y0, 5.9, 0.28,
                         font_size=10, color=DARK_GREY)
            y0 += 0.28
        y0 += 0.06

    # ── SLIDE 10 — SHAPE FEATURES & ENHANCEMENT TECHNIQUES ───────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Shape Features & Enhancement Techniques",
               "Morphology descriptors + preprocessing pipeline for field images")
    footer_bar(sl, "10 / 14")

    # Shape column
    add_rect(sl, 0.25, 1.35, 5.8, 0.4, ACCENT)
    add_text_box(sl, "📐  SHAPE FEATURES", 0.35, 1.38, 5.6, 0.35,
                 font_size=14, bold=True, color=WHITE)
    shape_rows = [
        ("Hu Moments  (7 invariants)",
         "Rotation, scale & translation invariant.\n"
         "Elongated (Sugar cane 8:1) vs compact (Bitter Leaf 1.5:1)."),
        ("Fourier Descriptors",
         "DFT of leaf boundary. High-freq. energy = serrations.\n"
         "False Daisy's deep serrations absent in smooth-margined herbs."),
        ("Geometric Properties",
         "Solidity: Momordica lobed ≈0.6 vs ovate ≈0.95\n"
         "Aspect ratio, circularity, eccentricity — fast & interpretable."),
        ("Deep Shape (YOLOv8 masks)",
         "32 prototype masks decompose instance boundary directly.\n"
         "Implicitly learns shape+texture jointly from training data."),
    ]
    y0 = 1.85
    for title, body in shape_rows:
        add_rect(sl, 0.25, y0, 5.8, 0.3, rgb(0xFF, 0xF8, 0xE1))
        add_text_box(sl, title, 0.35, y0+0.02, 5.6, 0.27,
                     font_size=11, bold=True, color=rgb(0x7A, 0x50, 0x00))
        y0 += 0.3
        for line in body.split("\n"):
            add_text_box(sl, f"  {line}", 0.35, y0, 5.6, 0.28,
                         font_size=10, color=DARK_GREY)
            y0 += 0.28
        y0 += 0.06

    # Enhancement column
    add_rect(sl, 6.8, 1.35, 6.28, 0.4, GREEN_MID)
    add_text_box(sl, "✨  ENHANCEMENT PIPELINE", 6.9, 1.38, 6.1, 0.35,
                 font_size=14, bold=True, color=WHITE)

    enhancements = [
        ("1. CLAHE",         "Tile-wise histogram equalisation (8×8 grid, clip=2.0).\n"
                             "Applied to L* channel in LAB. ★ Highest priority."),
        ("2. Bilateral Filter", "Edge-preserving smooth (σ_space=5, σ_color=30).\n"
                             "Keeps veins & serrations; removes sensor noise."),
        ("3. Multi-Scale Retinex", "3 scales σ={15,80,250}: removes illumination cast.\n"
                             "Critical for forest canopy vs direct sun variation."),
        ("4. Gamma Correction",  "I^(1/0.45) — recovers shadow detail.\n"
                             "Lightweight Retinex alternative for mobile devices."),
        ("5. Unsharp Masking",  "I + 0.5(I − blur). Sharpens venation & serrations.\n"
                             "Apply after bilateral, before LBP/Gabor extraction."),
    ]
    y0 = 1.85
    for title, body in enhancements:
        add_rect(sl, 6.8, y0, 6.28, 0.3, rgb(0xE8, 0xF5, 0xE9))
        add_text_box(sl, title, 6.9, y0+0.02, 6.1, 0.27,
                     font_size=11, bold=True, color=GREEN_DARK)
        y0 += 0.3
        for line in body.split("\n"):
            add_text_box(sl, f"  {line}", 6.9, y0, 6.1, 0.28,
                         font_size=10, color=DARK_GREY)
            y0 += 0.28
        y0 += 0.06

    add_text_box(sl,
        "Recommended order:  Raw image  →  CLAHE  →  Bilateral  →  Retinex  →  Unsharp  →  Feature extraction",
        0.25, 6.85, 12.83, 0.38,
        font_size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    add_rect(sl, 0.25, 6.82, 12.83, 0.42, GREEN_DARK)
    add_text_box(sl,
        "Recommended order:  Raw image  →  CLAHE  →  Bilateral  →  Retinex  →  Unsharp  →  Feature extraction",
        0.3, 6.85, 12.73, 0.38,
        font_size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

    # ── SLIDE 11 — RESULTS ────────────────────────────────────────────────────
    # (was slide 9 — number updated)
    results_pending = not METRICS_JSON.exists()
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Experimental Results", "Validation set performance")
    footer_bar(sl, "11 / 14")

    if results_pending:
        add_rect(sl, 1.5, 2.0, 10.33, 3.5, rgb(0xFF, 0xF8, 0xE1))
        add_rect(sl, 1.5, 2.0, 10.33, 0.6, ACCENT)
        add_text_box(sl, "Training In Progress — Results Pending",
                     1.6, 2.07, 10.13, 0.48,
                     font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl,
            "Once training completes, run:\n\n"
            "    python evaluate.py\n"
            "    python generate_presentation.py\n\n"
            "This slide will auto-populate with live metrics and per-class table.",
            1.7, 2.75, 9.93, 2.6,
            font_size=14, color=DARK_GREY, align=PP_ALIGN.LEFT)
        add_text_box(sl,
            "Preliminary estimate: mAP50 ~ 0.47  |  Best: False Daisy (0.995)  |  "
            "Challenging: Sugar cane (0.051)",
            0.5, 6.0, 12.33, 0.55,
            font_size=12, italic=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    else:
        mets = [
            (f"{ov['mAP50']:.2f}",    "mAP50\n(Box)",    GREEN_DARK),
            (f"{ov['mAP50_95']:.2f}", "mAP50-95\n(Box)", GREEN_MID),
            (f"{ov['precision']:.2f}","Precision",        GREEN_DARK),
            (f"{ov['recall']:.2f}",   "Recall",           GREEN_MID),
            (f"{ov['F1']:.2f}",       "F1 Score",         rgb(0x33, 0x69, 0x1E)),
        ]
        for i, (val, label, bg) in enumerate(mets):
            x = 0.3 + i * 2.55
            metric_box(sl, x, 1.4, 2.3, 1.5, val, label, bg)

        add_text_box(sl, "Per-Class Performance (Validation Set):",
                     0.3, 3.1, 12.5, 0.42,
                     font_size=14, bold=True, color=GREEN_DARK)

        headers = ["Class", "Precision", "Recall", "F1", "mAP50"]
        top_cls = sorted(per_class, key=lambda c: c["mAP50"], reverse=True)[:7]
        mp_   = sum(c["P"]     for c in per_class) / max(len(per_class), 1)
        mr_   = sum(c["R"]     for c in per_class) / max(len(per_class), 1)
        mf1_  = sum(c["F1"]    for c in per_class) / max(len(per_class), 1)
        map_  = sum(c["mAP50"] for c in per_class) / max(len(per_class), 1)
        rows = [
            (c["class"][:22], f"{c['P']:.2f}", f"{c['R']:.3f}", f"{c['F1']:.3f}", f"{c['mAP50']:.3f}")
            for c in top_cls
        ] + [(f"Mean ({len(per_class)} cls)", f"{mp_:.2f}", f"{mr_:.3f}", f"{mf1_:.3f}", f"{map_:.3f}")]

        col_xs = [0.3, 4.8, 6.8, 8.8, 10.8]
        col_ws = [4.3, 1.9, 1.9, 1.9, 2.3]

        y_h = 3.6
        for i, hdr in enumerate(headers):
            add_rect(sl, col_xs[i], y_h, col_ws[i], 0.38, GREEN_DARK)
            add_text_box(sl, hdr, col_xs[i]+0.05, y_h+0.04, col_ws[i]-0.05, 0.33,
                         font_size=12, bold=True, color=WHITE)

        for row_i, row in enumerate(rows):
            y_r = 3.98 + row_i * 0.34
            bg = LIGHT_GREY if row_i % 2 == 0 else WHITE
            if row_i == len(rows) - 1:
                bg = rgb(0xD7, 0xEB, 0xD7)
            for col_i, val in enumerate(row):
                add_rect(sl, col_xs[col_i], y_r, col_ws[col_i], 0.34, bg)
                add_text_box(sl, val,
                             col_xs[col_i]+0.05, y_r+0.03,
                             col_ws[col_i]-0.05, 0.30,
                             font_size=11, color=DARK_GREY,
                             align=PP_ALIGN.LEFT if col_i == 0 else PP_ALIGN.CENTER)

    # ── SLIDE 12 — CLASS HIGHLIGHTS ───────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Class Highlights", "Best and worst performing species")
    footer_bar(sl, "12 / 14")

    add_rect(sl, 0.3, 1.4, 5.9, 0.42, GREEN_MID)
    add_text_box(sl, "✅ Best Performing Classes", 0.4, 1.43, 5.8, 0.35,
                 font_size=14, bold=True, color=WHITE)

    # top 3 by mAP50 — dynamic
    top3 = sorted_cls[:3]
    for i, c in enumerate(top3):
        cls   = c["class"]
        score = f"mAP50 = {c['mAP50']:.3f}"
        reason = "High confidence — visually distinctive species"
        y = 1.95 + i * 0.8
        add_rect(sl, 0.3, y, 5.9, 0.72, LIGHT_GREY)
        add_text_box(sl, f"{cls}  |  {score}", 0.4, y+0.04, 5.7, 0.32,
                     font_size=12, bold=True, color=GREEN_DARK)
        add_text_box(sl, reason, 0.4, y+0.36, 5.7, 0.32,
                     font_size=11, color=DARK_GREY, italic=True)

    add_rect(sl, 6.8, 1.4, 6.2, 0.42, rgb(0x8B, 0x1A, 0x1A))
    add_text_box(sl, "⚠️  Challenging Classes", 6.9, 1.43, 6.1, 0.35,
                 font_size=14, bold=True, color=WHITE)

    # bottom 3 by mAP50 — dynamic
    bot3 = sorted_cls[-3:][::-1]  # worst first
    for i, c in enumerate(bot3):
        cls   = c["class"]
        score = f"mAP50 = {c['mAP50']:.3f}"
        reason = "Low performance — limited data or visually ambiguous"
        y = 1.95 + i * 0.8
        add_rect(sl, 6.8, y, 6.2, 0.72, rgb(0xFF, 0xF0, 0xF0))
        add_text_box(sl, f"{cls}  |  {score}", 6.9, y+0.04, 6.0, 0.32,
                     font_size=12, bold=True, color=rgb(0x8B, 0x1A, 0x1A))
        add_text_box(sl, reason, 6.9, y+0.36, 6.0, 0.32,
                     font_size=11, color=DARK_GREY, italic=True)

    add_text_box(sl, "Improvement path:", 0.3, 4.65, 12.5, 0.42,
                 font_size=14, bold=True, color=GREEN_DARK)
    for i, item in enumerate([
        "Collect 30–50 more images per weak class  |  Currently: avg 4.6 images/class",
        "Re-annotate 22 bbox images with accurate polygon masks",
        "Apply class-balanced sampling during training",
    ]):
        add_text_box(sl, f"• {item}", 0.4, 5.1 + i * 0.4, 12.4, 0.38,
                     font_size=12, color=DARK_GREY)

    # ── SLIDE 11 — CONTINUOUS LEARNING ────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
    header_bar(sl, "Continuous Learning",
               "Auto-triggered incremental fine-tuning")
    footer_bar(sl, "13 / 14")

    # Flow diagram
    steps = [
        ("1\nUser sees\nwrong label",    GREEN_DARK),
        ("2\nSelect correct\nherb species", GREEN_MID),
        ("3\nCorrection saved\nas YOLO label", GREEN_DARK),
        ("4\nDaemon thread\nlaunched", rgb(0x4A, 0x6A, 0x1A)),
        ("5\nFine-tuning\nruns in BG", ACCENT),
        ("6\nModel hot-\nreloaded", GREEN_DARK),
    ]
    for i, (label, bg) in enumerate(steps):
        x = 0.25 + i * 2.15
        add_rect(sl, x, 1.5, 1.95, 1.5, bg)
        add_text_box(sl, label, x, 1.55, 1.95, 1.4,
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(sl, "→", x+1.95, 2.0, 0.25, 0.6,
                         font_size=18, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

    details = [
        ("Low Learning Rate", "lr₀=0.0005\n(1/20 of initial)", GREEN_DARK),
        ("Epochs", "40 epochs\npatience=15", GREEN_MID),
        ("Start Point", "From best.pt\n(not scratch)", GREEN_DARK),
        ("Thread Safety", "threading.Lock\nprevents parallel runs", GREEN_MID),
    ]
    for i, (title, val, bg) in enumerate(details):
        x = 0.3 + i * 3.2
        add_rect(sl, x, 3.3, 2.9, 0.38, bg)
        add_text_box(sl, title, x+0.1, 3.32, 2.7, 0.32,
                     font_size=12, bold=True, color=WHITE)
        add_rect(sl, x, 3.68, 2.9, 0.85, LIGHT_GREY)
        add_text_box(sl, val, x+0.1, 3.73, 2.7, 0.75,
                     font_size=13, color=GREEN_DARK, align=PP_ALIGN.CENTER)

    add_rect(sl, 0.3, 4.75, 12.5, 0.05, GREEN_LIGHT)
    add_text_box(sl, "Why low LR prevents catastrophic forgetting:",
                 0.3, 4.85, 12.5, 0.42,
                 font_size=14, bold=True, color=GREEN_DARK)
    add_text_box(sl,
        "Small gradient steps on correction data barely move weights from the "
        "general herb-recognition manifold learned during full training. "
        "The model incorporates the specific correction without unlearning "
        "the 307 original annotations.",
        0.3, 5.3, 12.5, 0.8, font_size=12, color=DARK_GREY)

    add_text_box(sl,
        "Future work: Elastic Weight Consolidation (EWC) to formally measure and bound forgetting",
        0.3, 6.2, 12.5, 0.42, font_size=12, italic=True,
        color=rgb(0x5D, 0x40, 0x00))

    # ── SLIDE 12 — CONCLUSIONS & GITHUB ────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 13.33, 7.5, GREEN_DARK)
    add_rect(sl, 0, 5.8, 13.33, 1.7, rgb(0x10, 0x40, 0x10))
    add_rect(sl, 0, 1.0, 13.33, 0.05, ACCENT)

    add_text_box(sl, "Conclusions", 0.5, 0.1, 12.0, 0.8,
                 font_size=34, bold=True, color=WHITE)

    if results_pending:
        conc1 = "HerbScan: feature-aware YOLOv8n-seg for 16 East African medicinal herbs (results pending)"
        conc2 = "Colour (HSV/LAB/VegIdx), texture (Gabor/LBP/GLCM) and shape (Hu/Fourier) features guide recognition"
    else:
        conc1 = f"HerbScan achieves {ov['mAP50']:.2f} mAP50 on a 16-class herbal dataset"
        conc2 = (f"{best_cls['class']} & {second_best['class']} reach top accuracy "
                 f"(mAP50 {best_cls['mAP50']:.3f} & {second_best['mAP50']:.3f})")
    concs = [
        conc1,
        conc2,
        "Bbox→polygon conversion retains all 307 annotations — no data discarded",
        "Gradio 5 web interface supports camera capture, expert escalation, and feedback correction",
        "Background fine-tuning with low lr prevents catastrophic forgetting",
        "Deployed as open-source — all code available on GitHub",
    ]
    for i, c in enumerate(concs):
        y = 1.15 + i * 0.55
        add_rect(sl, 0.4, y, 0.38, 0.38, ACCENT)
        add_text_box(sl, str(i+1), 0.4, y+0.02, 0.38, 0.35,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, c, 0.85, y+0.03, 12.0, 0.45,
                     font_size=13, color=GREEN_LIGHT)

    add_rect(sl, 0.5, 4.8, 0.08, 0.9, ACCENT)
    add_text_box(sl, "GitHub Repository:", 0.7, 4.8, 12.0, 0.42,
                 font_size=16, bold=True, color=ACCENT)
    add_text_box(sl,
        "https://github.com/bertonag/HerbalSpiciesRecognition",
        0.7, 5.2, 12.0, 0.55, font_size=20, bold=True, color=WHITE)
    add_text_box(sl, "Clone, train, and run: python train.py → python app.py",
                 0.7, 5.75, 12.0, 0.42, font_size=13, italic=True, color=GREEN_LIGHT)

    add_text_box(sl, "Gilbert Nyakana  ·  Makerere University  ·  2026",
                 0.5, 6.3, 12.3, 0.4, font_size=13, color=GREEN_LIGHT,
                 align=PP_ALIGN.CENTER)
    add_text_box(sl, "14 / 14", 0.5, 6.9, 12.3, 0.4,
                 font_size=11, color=ACCENT, align=PP_ALIGN.RIGHT)

    prs.save(str(OUT))
    print(f"[OK] Presentation saved: {OUT}")


if __name__ == "__main__":
    build()
